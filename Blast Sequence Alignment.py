import mysql.connector
import subprocess
import vcf
import numpy as np
import matplotlib.pyplot as plt
from pptx import Presentation
from sklearn.ensemble import RandomForestClassifier

# MySQL Database Setup
def connect_db():
    db_connection = mysql.connector.connect(
        host="localhost",
        user="your_username",
        password="your_password",
        database="genomic_data"
    )
    return db_connection

# Insert sequence data into MySQL
def insert_sequence(name, sequence):
    db_connection = connect_db()
    cursor = db_connection.cursor()
    cursor.execute("""
    CREATE TABLE IF NOT EXISTS sequences (
        id INT AUTO_INCREMENT PRIMARY KEY,
        name VARCHAR(255),
        sequence TEXT
    )
    """)
    cursor.execute("INSERT INTO sequences (name, sequence) VALUES (%s, %s)", (name, sequence))
    db_connection.commit()
    cursor.close()
    db_connection.close()
    print(f"Inserted sequence: {name}")

# BLAST Sequence Alignment
def run_blast(query_file, database_file, output_file):
    command = f"blastn -query {query_file} -db {database_file} -out {output_file} -outfmt 6"
    subprocess.run(command, shell=True)
    print(f"BLAST analysis complete. Results saved to {output_file}")

# Parse BLAST results
def parse_blast_results(blast_file):
    with open(blast_file, 'r') as file:
        for line in file:
            columns = line.strip().split('\t')
            query, subject, identity, alignment_length = columns[0], columns[1], columns[2], columns[3]
            print(f"Query: {query}, Subject: {subject}, Identity: {identity}%")

# Analyze VCF for genetic variants
def analyze_variants(vcf_file):
    vcf_reader = vcf.Reader(filename=vcf_file)
    for record in vcf_reader:
        print(f"Variant: {record.CHROM}, {record.POS}, {record.ID}")

# Predict genetic variation outcomes using Random Forest
def predict_genetic_variation(X_train, y_train, X_new):
    model = RandomForestClassifier()
    model.fit(X_train, y_train)
    prediction = model.predict(X_new)
    return prediction

# Visualize Genetic Variation Frequencies
def visualize_variations():
    variations = [10, 20, 15, 30]
    labels = ['SNP1', 'SNP2', 'SNP3', 'SNP4']
    plt.bar(labels, variations)
    plt.xlabel('SNPs')
    plt.ylabel('Frequency')
    plt.title('Genetic Variation Frequencies')
    plt.show()

# Create a PowerPoint presentation
def create_presentation():
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Genetic Variation Insights"
    subtitle.text = "Presented to Farmers for Crop Improvement"

    slide_layout = prs.slide_layouts[5]  # Content slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Genetic Variation Frequencies"
    slide.shapes.add_picture("genetic_variation_plot.png", 0, 0, height=500, width=500)

    prs.save('Genetic_Variation_Presentation.pptx')
    print("Presentation created successfully.")

# Main function
def main():
    # Step 1: Insert sequence into database
    insert_sequence("Gene_A", "ATGCGTACGTTAG")
    insert_sequence("Gene_B", "ATCGTACGATGGC")

    # Step 2: Run BLAST and parse results
    run_blast("query_sequence.fasta", "genomic_database.fasta", "blast_results.txt")
    parse_blast_results("blast_results.txt")

    # Step 3: Analyze VCF file for genetic variants
    analyze_variants("variants_data.vcf")

    # Step 4: Train model and predict genetic variation
    X_train = np.array([[0, 1], [1, 0], [0, 0], [1, 1]])  # Example features
    y_train = np.array([0, 1, 0, 1])  # Labels (0: resistant, 1: not resistant)
    X_new = np.array([[1, 0]])  # New data for prediction
    prediction = predict_genetic_variation(X_train, y_train, X_new)
    print(f"Prediction for new genetic data: {prediction}")

    # Step 5: Visualize genetic variation frequencies
    visualize_variations()

    # Step 6: Create and save PowerPoint presentation
    create_presentation()

# Run main
if __name__ == "__main__":
    main()
